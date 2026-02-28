import os
import io
import json
import streamlit as st
import PyPDF2
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv

# Google SDK for Imagen & Text
from google import genai
from google.genai import types

# CrewAI & Tools
from crewai import Agent, Task, Crew, Process
from crewai_tools import DirectoryReadTool, FileReadTool, FileWriterTool
from crewai.tools import tool

# 환경 변수 로드
load_dotenv()

# ==========================================
# 1. 파일 텍스트 추출 (NotebookLM RAG 컴포넌트)
# ==========================================
def extract_text_from_file(uploaded_file):
    text = ""
    file_type = uploaded_file.name.split('.')[-1].lower()
    try:
        if file_type == 'pdf':
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
        elif file_type == 'pptx':
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif file_type == 'txt':
            text = uploaded_file.getvalue().decode("utf-8")
    except Exception as e:
        st.error(f"파일 텍스트 추출 중 오류 발생 ({uploaded_file.name}): {e}")
    return text

# ==========================================
# 2. 나노바나나 + PPTX 융합 도구 (Agent 용)
# ==========================================
@tool("CreatePptxTool")
def create_pptx_tool(json_slides_data: str) -> str:
    """
    JSON 슬라이드 데이터를 받아 10장 분량의 파워포인트(PPTX)를 생성합니다.
    데이터 포맷 예:
    [
      {"title": "...", "content": "...", "image_prompt": "상징적 비즈니스 이미지 프롬프트(영어)"}
    ]
    만약 image_prompt 가 있으면 나노바나나(Gemini Imagen)를 통해 이미지를 즉시 생성하고 삽입합니다.
    """
    try:
        slides_data = json.loads(json_slides_data)
        prs = Presentation()
        output_dir = os.environ.get("CURRENT_PROJECT_OUTPUT_DIR", "outputs")
        
        # 구글 API 키
        api_key = os.environ.get("GEMINI_API_KEY")

        for i, slide_info in enumerate(slides_data):
            if i >= 10: break
            title_text = slide_info.get("title", "")
            content_text = slide_info.get("content", "")
            image_prompt = slide_info.get("image_prompt", "")
            
            # 인메모리 이미지 스트림 폴백
            image_stream = None
            if image_prompt and api_key:
                try:
                    # NanoBanana (Imagen 3.0) 호출
                    client = genai.Client(api_key=api_key)
                    img_result = client.models.generate_images(
                        model='imagen-3.0-generate-001',
                        prompt=image_prompt,
                        config=types.GenerateImagesConfig(
                            number_of_images=1,
                            output_mime_type="image/jpeg",
                            aspect_ratio="4:3"
                        )
                    )
                    image_bytes = img_result.generated_images[0].image.image_bytes
                    image_stream = io.BytesIO(image_bytes)
                except Exception as img_err:
                    print(f"이미지 생성 에러 (슬라이드 {i+1}): {img_err}")
            
            # 슬라이드 생성 로직 (0번은 제목, 1번부터는 내용 템플릿)
            if i == 0:
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = title_text
                subtitle.text = content_text
                
                # 폰트 강제 주입
                for paragraph in title.text_frame.paragraphs:
                    for run in paragraph.runs: run.font.name = 'Malgun Gothic'
                for paragraph in subtitle.text_frame.paragraphs:
                    for run in paragraph.runs: run.font.name = 'Malgun Gothic'
                    
                if image_stream:
                    # 제목 슬라이드의 경우 중앙 하단 등에 배치 (임의 조절 가능)
                    slide.shapes.add_picture(image_stream, Inches(3.5), Inches(4.5), height=Inches(2.5))
            else:
                slide_layout = prs.slide_layouts[1] # 빈 레이아웃(6) 등을 쓰면 박스 조정이 더 편함
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                body = slide.placeholders[1]
                
                title.text = title_text
                body.text = content_text
                
                # 우측 영역에 이미지를 삽입하기 위해 기본 텍스트 박스의 너비를 줄임
                if image_stream:
                    body.width = Inches(4.5)
                    # 우측 상단/중앙 위치에 이미지 삽입
                    slide.shapes.add_picture(image_stream, Inches(5.0), Inches(1.5), height=Inches(4.5))

                # 폰트 강제 주입
                for paragraph in title.text_frame.paragraphs:
                    for run in paragraph.runs: run.font.name = 'Malgun Gothic'
                if body.has_text_frame:
                    for paragraph in body.text_frame.paragraphs:
                        for run in paragraph.runs: run.font.name = 'Malgun Gothic'
                        
        output_path = os.path.join(output_dir, "presentation.pptx")
        prs.save(output_path)
        return f"나노바나나 연동 성공: {output_path} 파일이 자동 렌더링되어 생성되었습니다."
    except Exception as e:
        return f"PPTX 파일 생성 중 치명적 오류가 발생했습니다: {str(e)}"

# ==========================================
# 3. CrewAI 에이전트 팩토리 함수
# ==========================================
def run_crew(project_name, topic, context_text, llm_configs, task_flags):
    base_dir = os.path.join("projects", project_name)
    output_dir = os.path.join(base_dir, "outputs")
    os.environ["CURRENT_PROJECT_OUTPUT_DIR"] = output_dir

    run_design = task_flags.get("run_design", True)
    run_code = task_flags.get("run_code", True)
    run_ppt = task_flags.get("run_ppt", True)

    file_writer_tool = FileWriterTool()
    file_read_tool = FileReadTool()
    directory_read_tool = DirectoryReadTool(directory=output_dir)

    agent_tools = [directory_read_tool, file_read_tool, file_writer_tool]
    ppt_tools = [directory_read_tool, file_read_tool, file_writer_tool, create_pptx_tool]

    # [에이전트 정의]
    chief_pm = Agent(
        role='카리스로 수석 PM',
        goal=f'"{topic}" 프로젝트의 상세 기획안(PRD)을 깊이 있게 작성하고 {output_dir}/plan.md 파일에 저장함',
        backstory='당신은 NotebookLM 데이터를 최우선적으로 신뢰하는 카리스로의 리더입니다. 제공된 자료의 인사이트를 철저하게 기획과 구조에 반영하며, 문서 작성 능력이 탁월합니다. 항상 한국어로 작성합니다.',
        llm=llm_configs['pm_llm'],
        verbose=True,
        allow_delegation=True,
        tools=agent_tools
    )

    designer = Agent(
        role='UX/UI 전략가',
        goal=f'기획안을 바탕으로 Tailwind CSS 스타일 시트와 UI 컴포넌트 구조를 작성하고 {output_dir}/style.css 및 {output_dir}/ui_components.md로 저장함',
        backstory='미니멀리즘을 선호하는 베테랑 디자이너입니다. Tailwind CSS에 능통합니다.',
        llm=llm_configs['designer_llm'],
        verbose=True,
        tools=agent_tools
    )

    developer = Agent(
        role='수석 엔지니어',
        goal=f'기획안과 디자인을 참조하여 Next.js/Node.js 핵심 로직 코드를 구상하고 {output_dir}/app.js에 저장함',
        backstory='확장성 있는 코드를 작성하는 풀스택 개발 전문가입니다. 최신 React 생태계와 Node.js 백엔드를 설계합니다.',
        llm=llm_configs['developer_llm'],
        verbose=True,
        tools=agent_tools
    )

    ppt_master = Agent(
        role='PPT 마스터',
        goal=f'결과물을 취합해 발표 대본({output_dir}/pitch_script.md)을 정리하고, 나노바나나용 Image Prompt를 담은 JSON을 도출해 파워포인트 파일(presentation.pptx)을 자동 생성함',
        backstory='비즈니스 문서를 시각적으로 완벽히 정리하는 전문가. 깊이 있는 자료를 바탕으로 상징적이고 전문적인 비즈니스 이미지 프롬프트(반드시 영문)를 함께 설계해내는 나노바나나 조련사입니다.',
        llm=llm_configs['ppt_llm'],
        verbose=True,
        tools=ppt_tools
    )

    # [핵심 Task 정의 - 동적 연결]
    plan_description = f'''
    1. 주제: "{topic}"
    2. [NotebookLM Primary Context - 절대적 우선순위 자료]: 
    아래 제공된 참고 문서(Context)가 있다면 다른 지식보다 가장 최우선적으로 존중하고 모든 기획과 목차 설정, 방향성에 그대로 반영하세요.
    [참고 문서 내용 시작]
    {context_text}
    [참고 문서 내용 끝]
    3. 반드시 `FileWriterTool`을 사용하여 한국어 베이스의 상세 PRD 작업 결과를 `{output_dir}/plan.md` 파일로 저장하세요.
    '''
    
    task_plan = Task(
        description=plan_description,
        expected_output=f"NotebookLM 데이터가 면밀하게 반영된 상세 PRD가 포함된 {output_dir}/plan.md 파일",
        agent=chief_pm
    )

    active_tasks = [task_plan]
    active_agents = [chief_pm]

    if run_design:
        task_design = Task(
            description=f'''
            1. `FileReadTool`을 사용하여 `{output_dir}/plan.md` 파일을 읽어 프로젝트 기획안을 파악하세요.
            2. 기획안을 바탕으로 매력적인 Tailwind CSS 기반의 스타일 시트를 작성해 `{output_dir}/style.css`로 저장하세요.
            3. UI 컴포넌트 구조 명세서를 작성해 `{output_dir}/ui_components.md`로 저장하세요.
            ''',
            expected_output=f"{output_dir}/style.css 및 {output_dir}/ui_components.md 파일",
            agent=designer,
            context=[active_tasks[-1]] 
        )
        active_tasks.append(task_design)
        active_agents.append(designer)

    if run_code:
        task_develop = Task(
            description=f'''
            1. `DirectoryReadTool` 및 `FileReadTool`을 사용해 `{output_dir}/plan.md` 등 앞선 파일들을 참조하세요.
            2. 실제 동작 가능한 Next.js/Node.js 핵심 로직 코드(API 및 컴포넌트 구조)를 작성하세요.
            3. 코드를 `{output_dir}/app.js`에 저장하세요.
            ''',
            expected_output=f"{output_dir}/app.js 파일",
            agent=developer,
            context=[active_tasks[-1]]
        )
        active_tasks.append(task_develop)
        active_agents.append(developer)

    if run_ppt:
        task_document = Task(
            description=f'''
            1. 앞선 기획안(plan.md) 및 파일들을 읽어 종합하여 전문적인 발표 대본을 작성하고 `{output_dir}/pitch_script.md`에 한국어로 저장하세요.
            2. 작성한 구조를 바탕으로 총 10장 구성의 슬라이드 데이터를 JSON 포맷으로 도출하세요. 
               JSON 객체 필수 속성: {{"title": "...", "content": "...", "image_prompt": "An abstract, hyper-realistic image of [concept] in corporate style..."}}
               *주의: image_prompt는 반드시 상세하고 시각적인 묘사를 담은 **영어(English)** 문장이어야 합니다.*
            3. 도출된 JSON 문자열을 `CreatePptxTool`에 넘겨 호출하여 `{output_dir}/presentation.pptx` 나노바나나 연동 파일 자동 생성을 완료하세요.
            ''',
            expected_output=f"{output_dir}/pitch_script.md 및 {output_dir}/presentation.pptx (나노바나나 첨부된 실제 파일 생성 증빙)",
            agent=ppt_master,
            context=[active_tasks[-1]]
        )
        active_tasks.append(task_document)
        active_agents.append(ppt_master)

    charisro_team = Crew(
        agents=active_agents,
        tasks=active_tasks,
        process=Process.sequential
    )
    
    return charisro_team.kickoff(inputs={'topic': topic})

# ==========================================
# 4. Streamlit UI (main)
# ==========================================
def main():
    st.set_page_config(page_title="Charisro AI Dashboard", page_icon="🤖", layout="wide")
    
    st.title("🤖 카리스로(Charisro) AI: 나노바나나 업그레이드 에디션")
    st.markdown("수석 PM, UX/UI 전문가, 개발자, PPT 마스터가 선택적으로 협업하며 일체형 결과물을 파일 시스템으로 도출합니다.")
    
    # 사이드바 설정 영역
    with st.sidebar:
        st.header("1. API Key 연동 상태")
        
        # 나노바나나(Gemini) 통합 키 프로세스
        gemini_key = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")
        if gemini_key:
            st.success("✅ Gemini (나노바나나) API Key 등록됨")
        else:
            new_gemini = st.text_input("Gemini API Key (플래닝 및 나노바나나 이미지 필수)", type="password")
            if new_gemini: os.environ["GEMINI_API_KEY"] = new_gemini
            
        openai_key = os.getenv("OPENAI_API_KEY")
        if openai_key:
            st.success("✅ OpenAI API Key 등록됨")
        else:
            new_openai = st.text_input("OpenAI API Key (미입력 시 에러 가능)", type="password")
            if new_openai: os.environ["OPENAI_API_KEY"] = new_openai
            
        anthropic_key = os.getenv("ANTHROPIC_API_KEY")
        if anthropic_key:
            st.success("✅ Anthropic API Key 등록됨")
        else:
            new_anthropic = st.text_input("Anthropic API Key (Claude 구동 시 필요)", type="password")
            if new_anthropic: os.environ["ANTHROPIC_API_KEY"] = new_anthropic

        st.divider()
        st.header("2. 에이전트 브레인 선택")
        pm_model = st.selectbox("수석 PM (기획)", ["gemini/gemini-2.5-flash", "gemini/gemini-2.5-pro", "gemini/gemini-1.5-pro"])
        designer_model = st.selectbox("UX/UI 전략가", ["openai/gpt-4o", "openai/gpt-4-turbo"])
        developer_model = st.selectbox("수석 엔지니어 (코드)", ["anthropic/claude-3-opus-20240229", "anthropic/claude-sonnet-4-6", "anthropic/claude-3-5-sonnet-20241022", "gemini/gemini-2.5-flash"])
        ppt_model = st.selectbox("PPT 마스터 (대본/나노바나나)", ["openai/gpt-4o", "openai/gpt-4-turbo"])
        
        st.divider()
        st.header("3. 프로젝트 폴더 설정")
        project_name = st.text_input("프로젝트 이름", value="demo_project")
        
    llm_configs = {
        'pm_llm': pm_model,
        'designer_llm': designer_model,
        'developer_llm': developer_model,
        'ppt_llm': ppt_model
    }

    base_dir = os.path.join("projects", project_name)
    uploads_dir = os.path.join(base_dir, "uploads")
    outputs_dir = os.path.join(base_dir, "outputs")
    
    # 메인 영역 상단 분리
    st.subheader("📚 1. 카리스로 작업 파이프라인 (원하는 역할만 지시하세요!)")
    col0_1, col0_2, col0_3, col0_4 = st.columns(4)
    with col0_1:
        st.info("🦸‍♂️ 수석 PM")
        run_plan = st.checkbox("문서 분석 및 상세 기획안 도출", value=True, disabled=True, help="기획은 전체 파이프라인의 핵심이므로 끌 수 없습니다.")
    with col0_2:
        st.info("👩‍🎨 디자이너")
        run_design = st.checkbox("디자인 및 UI (Tailwind) 명세 작성", value=True)
    with col0_3:
        st.info("👩‍💻 개발자")
        run_code = st.checkbox("Next.js / Node.js 핵심 코드 설계", value=False)
    with col0_4:
        st.info("📊 PPT 마스터")
        run_ppt = st.checkbox("나노바나나 이미지 연동 PPTX 자동생성", value=True)

    task_flags = {
        "run_design": run_design,
        "run_code": run_code,
        "run_ppt": run_ppt
    }

    col1, col2 = st.columns([2, 1])
    with col1:
        st.subheader("📝 2. 기획 지시사항")
        topic = st.text_area("프로젝트 주제, 방향성 등 상세한 지시를 내려주세요.", height=120)
        
    with col2:
        st.subheader("📁 3. NotebookLM 딥다운 컨텍스트")
        uploaded_files = st.file_uploader(
            "강력하게 참고해야 할 지식(PDF, PPTX, TXT)을 업로드하세요.",
            type=["pdf", "pptx", "txt"],
            accept_multiple_files=True
        )
        
    if st.button("🚀 카리스로 맞춤형 군단 가동하기", use_container_width=True):
        if not topic.strip():
            st.error("프로젝트 주제를 입력해주세요!")
            return
            
        used_providers = [m.split('/')[0] for m in llm_configs.values()]
        if "gemini" in used_providers and not os.getenv("GEMINI_API_KEY"):
            st.error("Gemini 모델이 선택되었지만 API Key가 누락되었습니다. 사이드바에 입력해주세요.")
            return
        if "openai" in used_providers and not os.getenv("OPENAI_API_KEY"):
            st.error("OpenAI 모델이 선택되었지만 API Key가 누락되었습니다. 사이드바에 입력해주세요.")
            return
        if run_code and "anthropic" in developer_model and not os.getenv("ANTHROPIC_API_KEY"):
            st.error("Anthropic(Claude) 모델이 코딩에 할당되었지만 API Key가 누락되었습니다.")
            return
            
        os.makedirs(uploads_dir, exist_ok=True)
        os.makedirs(outputs_dir, exist_ok=True)
        
        context_text = ""
        if uploaded_files:
            with st.spinner("NotebookLM 스타일 다중 문서 심층 분석 중..."):
                for uf in uploaded_files:
                    file_path = os.path.join(uploads_dir, uf.name)
                    with open(file_path, "wb") as f:
                        f.write(uf.getbuffer())
                    extracted = extract_text_from_file(uf)
                    if extracted:
                        context_text += f"\n\n--- [참고원문: {uf.name}] ---\n" + extracted
            st.success("데이터 추출 및 Context 주입 완료! 수석 PM이 이를 최우선으로 참고합니다.")
            
        with st.status("🧠 에이전트들이 선택된 파이프라인만 집중적으로 수행 중입니다...", expanded=True) as status:
            try:
                result = run_crew(project_name, topic, context_text, llm_configs, task_flags)
                status.update(label="✅ 맞춤형 작업이 성공적으로 완료되었습니다!", state="complete", expanded=False)
            except Exception as e:
                status.update(label="❌ 에러가 발생했습니다.", state="error")
                st.error(f"실행 중 구조적 오류가 발생했습니다: {str(e)}")
                return
                
        st.success("모든 처리가 끝났습니다. 하단 대시보드에서 산출물을 확인하세요.")
        
    st.divider()
    
    st.header("📊 커스텀 결과물 대시보드")
    if os.path.exists(outputs_dir):
        files = os.listdir(outputs_dir)
        if files:
            tabs_data = {
                "기획서 (plan.md)": "plan.md",
                "디자인 (style.css, ui.md)": None, # Complex
                "코드 (app.js)": "app.js",
                "대본 (pitch_script.md)": "pitch_script.md",
                "발표파일 (presentation.pptx)": "presentation.pptx",
            }
            
            tabs = st.tabs(list(tabs_data.keys()))
            
            def read_file(filename):
                path = os.path.join(outputs_dir, filename)
                if os.path.exists(path):
                    with open(path, "r", encoding="utf-8") as f:
                        return f.read()
                return "현재 사이클에서 생략되었거나, 아직 산출되지 않은 결과물입니다."
                
            def get_file_bytes(filename):
                path = os.path.join(outputs_dir, filename)
                if os.path.exists(path):
                    with open(path, "rb") as f:
                        return f.read()
                return None
            
            with tabs[0]:
                st.markdown(read_file("plan.md"))
                bt = get_file_bytes("plan.md")
                if bt: st.download_button("📝 plan.md 다운로드", data=bt, file_name="plan.md", mime="text/markdown")
            
            with tabs[1]:
                ui_content = read_file("ui_components.md")
                css_content = read_file("style.css")
                st.subheader("UI 명세")
                st.markdown(ui_content)
                bt1 = get_file_bytes("ui_components.md")
                if bt1: st.download_button("🧩 ui_components.md 다운로드", data=bt1, file_name="ui_components.md")
                
                st.subheader("스타일 시트")
                st.code(css_content, language="css")
                bt2 = get_file_bytes("style.css")
                if bt2: st.download_button("🎨 style.css 다운로드", data=bt2, file_name="style.css")
                
            with tabs[2]:
                st.code(read_file("app.js"), language="javascript")
                bt = get_file_bytes("app.js")
                if bt: st.download_button("💻 app.js 다운로드", data=bt, file_name="app.js")
                
            with tabs[3]:
                st.markdown(read_file("pitch_script.md"))
                bt = get_file_bytes("pitch_script.md")
                if bt: st.download_button("🎙️ pitch_script.md 다운로드", data=bt, file_name="pitch_script.md")
                
            with tabs[4]:
                st.info("나노바나나 협동: PPT 마스터가 AI 상징 레이아웃을 접목하여 자동 생성한 발표 파일입니다.")
                bt = get_file_bytes("presentation.pptx")
                if bt: 
                    st.download_button(
                        label="📥 나노바나나 presentation.pptx 다운로드", 
                        data=bt, file_name="presentation.pptx", 
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        type="primary"
                    )
        else:
            st.info("아직 생성된 결과물이 없습니다.")
    else:
        st.info("프로젝트 폴더가 비어있습니다.")

if __name__ == "__main__":
    main()
