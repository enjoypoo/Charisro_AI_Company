import os
import json
from dotenv import load_dotenv
from crewai import Agent, Task, Crew, Process
from crewai_tools import DirectoryReadTool, FileReadTool, FileWriterTool
from crewai.tools import tool
from pptx import Presentation
from pptx.util import Inches, Pt

# 환경 변수 로드
load_dotenv()

@tool("CreatePptxTool")
def create_pptx_tool(json_slides_data: str) -> str:
    """
    이 도구는 JSON 형태의 슬라이드 데이터를 받아 10장 분량의 실제 파워포인트 파일(.pptx)을 생성하고 
    저장합니다. 데이터 형식은 다음과 같은 JSON 문자열이어야 합니다:
    [
      {"title": "슬라이드 1 제목", "content": "본문 내용이나 핵심 문구"},
      {"title": "슬라이드 2 제목", "content": "1. 포인트 일\n2. 포인트 이"}
    ]
    최대 10개의 슬라이드를 생성합니다.
    """
    try:
        slides_data = json.loads(json_slides_data)
        
        prs = Presentation()
        
        for i, slide_info in enumerate(slides_data):
            if i >= 10:  # 10장 제한
                break
                
            title_text = slide_info.get("title", "")
            content_text = slide_info.get("content", "")
            
            # 첫 번째 슬라이드는 제목 슬라이드 레이아웃(0)
            if i == 0:
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                title.text = title_text
                subtitle.text = content_text
                
                # 폰트 강제 적용 (한글 깨짐 방지)
                for paragraph in title.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Malgun Gothic'
                for paragraph in subtitle.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Malgun Gothic'
                        
            else:
                # 일반 슬라이드는 제목 및 내용 레이아웃(1)
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                body = slide.placeholders[1]
                
                title.text = title_text
                body.text = content_text
                
                # 폰트 강제 적용
                for paragraph in title.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Malgun Gothic'
                for paragraph in body.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Malgun Gothic'
                        
        output_path = os.path.join("outputs", "presentation.pptx")
        prs.save(output_path)
        return f"성공적으로 {output_path} 파일이 생성되었습니다."
    
    except Exception as e:
        return f"PPTX 파일 생성 중 오류가 발생했습니다: {str(e)}"


def run_crew(topic: str):
    print(f"\n[{topic}] 자율형 프로젝트를 위한 에이전트 군단을 초기화합니다...")
    
    # 작업물 디렉토리 생성
    output_dir = 'outputs'
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        print(f"[{output_dir}] 폴더를 생성했습니다.")

    # 도구 초기화
    file_writer_tool = FileWriterTool()
    file_read_tool = FileReadTool()
    directory_read_tool = DirectoryReadTool(directory=output_dir)

    # 기본 요원 도구
    agent_tools = [directory_read_tool, file_read_tool, file_writer_tool]
    
    # PPT 마스터용 도구 (VBA 대신 pptx 직접 생성기 포함)
    ppt_tools = [directory_read_tool, file_read_tool, file_writer_tool, create_pptx_tool]

    # [카리스로(Charisro) 에이전트 군단 정의]
    chief_pm = Agent(
        role='카리스로 수석 PM',
        goal=f'"{topic}" 프로젝트의 상세 기획안(PRD)을 작성하고 outputs/plan.md 파일에 저장함',
        backstory='당신은 카리스로의 리더입니다. 논리적이고 체계적인 의사결정 전문가이자 훌륭한 문서 작성자입니다. 항상 한국어로 작성하며 전문 기술 용어는 영어(English)와 병기합니다.',
        llm='gemini/gemini-2.5-flash',
        verbose=True,
        allow_delegation=True,
        tools=agent_tools
    )

    designer = Agent(
        role='UX/UI 전략가',
        goal='완성된 기획안(plan.md)을 바탕으로 Tailwind CSS 스타일 시트와 UI 컴포넌트 구조를 작성하고 outputs/style.css 및 outputs/ui_components.md에 저장함',
        backstory='미니멀리즘을 선호하는 베테랑 디자이너입니다. Tailwind CSS에 능통합니다. 항상 한국어로 작성하며 전문 기술 용어는 영어(English)와 병기합니다.',
        llm='openai/gpt-4o',
        verbose=True,
        tools=agent_tools
    )

    developer = Agent(
        role='수석 엔지니어',
        goal='기획안과 디자인 파일을 참조하여 실제 동작 가능한 Next.js/Node.js 핵심 로직 코드를 작성하고 outputs/app.js에 저장함',
        backstory='확장성 있는 코드를 작성하는 풀스택 개발 전문가입니다. 최신 React 생태계와 Node.js 백엔드를 설계합니다. 항상 한국어로 작성하며 전문 기술 용어는 영어(English)와 병기합니다.',
        llm='gemini/gemini-2.5-flash',
        verbose=True,
        tools=agent_tools
    )

    ppt_master = Agent(
        role='PPT 마스터',
        goal='모든 결과물을 취합해 투자 유치용 대본(pitch_script.md)을 작성하고, CreatePptxTool을 사용해 파워포인트 파일(presentation.pptx)을 직접 생성함',
        backstory='비즈니스 문서를 시각적으로 완벽하게 정리하고 발표 대본을 매끄럽게 작성하는 전문가입니다. 파이썬 기반의 PPTX 생성 도구를 활용할 줄 압니다. 항상 한국어로 작성하며 전문 기술 용어는 영어(English)와 병기합니다.',
        llm='openai/gpt-4o',
        verbose=True,
        tools=ppt_tools
    )

    # [업무(Task) 정의 - 체인 연결]
    
    task_plan = Task(
        description=f'''
        1. 주제: "{topic}"
        2. 시장 분석과 수익 모델(BM)을 포함한 상세 기획안(Product Requirements Document, PRD)을 작성하세요.
        3. 반드시 `FileWriterTool`을 사용하여 작업 결과를 `outputs/plan.md` 파일로 저장하세요.
        ''',
        expected_output="상세 PRD가 포함된 outputs/plan.md 파일",
        agent=chief_pm
    )

    task_design = Task(
        description='''
        1. `FileReadTool`을 사용하여 `outputs/plan.md` 파일을 읽어 프로젝트 기획안을 파악하세요.
        2. 기획안을 바탕으로 매력적인 Tailwind CSS 기반의 스타일 시트를 작성하고 `FileWriterTool`을 사용해 `outputs/style.css`로 저장하세요.
        3. UI 컴포넌트 구조 명세서를 작성하고 `FileWriterTool`을 사용해 `outputs/ui_components.md`로 저장하세요.
        ''',
        expected_output="작성 완료된 outputs/style.css 및 outputs/ui_components.md 파일",
        agent=designer,
        context=[task_plan]
    )

    task_develop = Task(
        description='''
        1. `DirectoryReadTool` 및 `FileReadTool`을 사용해 `outputs/plan.md`와 `outputs/ui_components.md`, `outputs/style.css`를 참조하세요.
        2. 기획과 디자인 가이드를 바탕으로 실제 동작 가능한 Next.js/Node.js 핵심 로직 코드(API 및 메인 컴포넌트 구조)를 작성하세요.
        3. `FileWriterTool`을 사용하여 코드를 `outputs/app.js` (또는 필요시 여러 파일)에 저장하세요.
        ''',
        expected_output="동작 가능한 구조가 담긴 outputs/app.js 파일",
        agent=developer,
        context=[task_design]
    )

    task_document = Task(
        description='''
        1. `DirectoryReadTool`로 outputs 디렉토리를 살피고, `FileReadTool`로 앞선 기획, 디자인, 개발 파일들을 읽으세요.
        2. 위 내용들을 종합하여 투자 유치용 발표 자료 대본을 작성하고 `FileWriterTool`을 사용해 `outputs/pitch_script.md`에 저장하세요.
        3. 작성한 대본을 바탕으로 총 10장 구성의 (제목, 목차, 본문 등) 슬라이드 데이터를 JSON 포맷으로 구성하세요. 각 슬라이드의 형태는 {{"title": "...", "content": "..."}} 여야 합니다.
        4. 구성한 JSON 문자열을 `CreatePptxTool`의 인자로 넘겨 호출하세요. 이렇게 하면 `outputs/presentation.pptx` 실제 파일이 자동 생성됩니다.
        ''',
        expected_output="최종 취합된 outputs/pitch_script.md 및 outputs/presentation.pptx (실제 파일 생성 증빙) 결과",
        agent=ppt_master,
        context=[task_develop]
    )

    # [크루(Crew) 설정 및 가동]
    charisro_team = Crew(
        agents=[chief_pm, designer, developer, ppt_master],
        tasks=[task_plan, task_design, task_develop, task_document],
        process=Process.sequential
    )
    
    # 시스템 실행
    return charisro_team.kickoff(inputs={'topic': topic})

if __name__ == "__main__":
    print("\n### 카리스로(Charisro) 자율형 AI 시스템 가동 ###")
    user_query = input("어떤 프로젝트를 시작할까요? : ")
    
    result = run_crew(user_query)
    
    print("\n\n###############################")
    print("## 카리스로 완전 자율형 프로젝트 종료 ##")
    print("###############################\n")
    print("outputs/ 폴더에 저장된 파일들을 확인해주세요.")
    print(result)